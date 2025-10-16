import os
import json
import boto3
from datetime import datetime
import tempfile
from pptx import Presentation
from pptx.util import Inches, Pt

def handler(event, context):
    try:
        print(f"[POWERPOINT] Received event: {json.dumps(event)}")
        
        # Check if this is a Bedrock Agent action group invocation
        if 'agent' in event or 'actionGroup' in event or 'function' in event:
            print(f"[POWERPOINT] Bedrock Agent action group invocation detected")
            
            parameters = event.get('parameters', [])
            param_dict = {}
            for param in parameters:
                param_name = param['name']
                param_value = param['value']
                
                # Handle nested JSON in parameters
                if param_name in ['proposal_data']:
                    try:
                        if isinstance(param_value, str):
                            param_dict[param_name] = json.loads(param_value)
                        else:
                            param_dict[param_name] = param_value
                    except json.JSONDecodeError as e:
                        print(f"[POWERPOINT] JSON decode error for {param_name}: {str(e)}")
                        param_dict[param_name] = param_value
                else:
                    param_dict[param_name] = param_value
            
            session_id = param_dict.get('session_id')
            template_path = param_dict.get('template_path')
            proposal_data = param_dict.get('proposal_data')
        elif 'inputText' in event:
            try:
                params = json.loads(event['inputText'])
                session_id = params['session_id']
                template_path = params.get('template_path')
                proposal_data = params.get('proposal_data')
                if isinstance(proposal_data, str):
                    proposal_data = json.loads(proposal_data)
            except (json.JSONDecodeError, KeyError) as e:
                return format_agent_response(400, f'Invalid input format from AgentCore: {str(e)}')
        else:
            return format_agent_response(400, 'Missing required parameters')

        print(f"[POWERPOINT] Generating presentation for session: {session_id}")
        
        # Initialize AWS clients
        s3 = boto3.client('s3')
        dynamodb = boto3.client('dynamodb')
        bedrock = boto3.client('bedrock-runtime')
        
        # Update status
        dynamodb.update_item(
            TableName=os.environ['SESSIONS_TABLE_NAME'],
            Key={'session_id': {'S': session_id}},
            UpdateExpression='SET #status = :status, updated_at = :ua',
            ExpressionAttributeNames={'#status': 'status'},
            ExpressionAttributeValues={
                ':status': {'S': 'GENERATING_POWERPOINT'},
                ':ua': {'S': datetime.utcnow().isoformat()}
            }
        )
        
        # Get complete session data including raw requirements
        response = dynamodb.get_item(
            TableName=os.environ['SESSIONS_TABLE_NAME'],
            Key={'session_id': {'S': session_id}}
        )
        
        if 'Item' not in response:
            raise ValueError(f"Session {session_id} not found")
        
        session_data = response['Item']
        requirements_data_str = session_data.get('requirements_data', {}).get('S', '{}')
        cost_data_str = session_data.get('cost_data', {}).get('S', '{}')
        
        try:
            requirements_data_full = json.loads(requirements_data_str)
            raw_requirements = requirements_data_full.get('raw_requirements', '')
            analysis_data = requirements_data_full.get('analysis', {})
            cost_data = json.loads(cost_data_str)
        except json.JSONDecodeError:
            raw_requirements = ''
            analysis_data = {}
            cost_data = {}
        
        # Use Bedrock to generate presentation content with full context
        prompt = f"""Generate content for a professional PowerPoint presentation based on the following project information.

RAW PROJECT REQUIREMENTS (Original User Input):
{raw_requirements}

ANALYZED PROJECT DATA:
{json.dumps(analysis_data, indent=2)}

COST DATA:
{json.dumps(cost_data, indent=2)}

Create slides for a compelling project proposal presentation. Return as JSON (no markdown, no code blocks):
{{
    "slides": [
        {{
            "title": "Slide title",
            "content": ["Bullet point 1", "Bullet point 2"],
            "notes": "Speaker notes"
        }}
    ]
}}

Include slides for: Executive Summary, Project Scope, Technical Approach, Timeline, Team & Resources, Cost Breakdown, Benefits, and Next Steps."""
        
        print(f"[POWERPOINT] Calling Bedrock for presentation content generation...")
        
        bedrock_response = bedrock.invoke_model(
            modelId=os.environ['BEDROCK_MODEL_ID'],
            contentType='application/json',
            body=json.dumps({
                "messages": [
                    {
                        "role": "user",
                        "content": [{"text": prompt}]
                    }
                ],
                "inferenceConfig": {
                    "max_new_tokens": 3000,
                    "temperature": 0.7
                }
            })
        )
        
        result = json.loads(bedrock_response['body'].read().decode())
        
        # Parse the response - handle different response structures
        ppt_content = None
        if 'content' in result and len(result['content']) > 0:
            ppt_content = result['content'][0].get('text', '')
        elif 'output' in result:
            if isinstance(result['output'], dict) and 'message' in result['output']:
                message = result['output']['message']
                if 'content' in message and len(message['content']) > 0:
                    ppt_content = message['content'][0].get('text', '')
        
        if not ppt_content:
            print(f"[POWERPOINT] WARNING: Could not extract text from model response")
            # Use fallback data
            slides_data = {
                'slides': [
                    {
                        'title': 'Project Overview',
                        'content': [analysis_data.get('project_scope', 'Project overview')],
                        'notes': ''
                    }
                ]
            }
        else:
            # Parse JSON from response
            try:
                if '```json' in ppt_content:
                    ppt_content = ppt_content.split('```json')[1].split('```')[0].strip()
                elif '```' in ppt_content:
                    ppt_content = ppt_content.split('```')[1].split('```')[0].strip()
                
                slides_data = json.loads(ppt_content)
            except json.JSONDecodeError as e:
                print(f"[POWERPOINT] Could not parse presentation content as JSON: {str(e)}")
                # Use fallback
                slides_data = {
                    'slides': [
                        {
                            'title': 'Project Overview',
                            'content': [analysis_data.get('project_scope', 'Project overview')],
                            'notes': ''
                        }
                    ]
                }
        
        # Download template or create new presentation
        prs = None
        if template_path and not template_path.startswith('path/to/'):
            # Check if template exists in S3 before downloading
            try:
                s3.head_object(
                    Bucket=os.environ['TEMPLATES_BUCKET_NAME'],
                    Key=template_path
                )
                # Template exists, download it
                with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as template_file:
                    s3.download_file(
                        os.environ['TEMPLATES_BUCKET_NAME'],
                        template_path,
                        template_file.name
                    )
                    prs = Presentation(template_file.name)
                    print(f"[POWERPOINT] Using template: {template_path}")
            except s3.exceptions.NoSuchKey:
                print(f"[POWERPOINT] Template not found: {template_path}, creating blank presentation")
                prs = Presentation()
            except Exception as e:
                print(f"[POWERPOINT] Error loading template: {str(e)}, creating blank presentation")
                prs = Presentation()
        else:
            print(f"[POWERPOINT] No valid template specified, creating blank presentation")
            prs = Presentation()
        
        # Add slides based on generated content
        for slide_data in slides_data.get('slides', []):
            slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title and Content layout
            
            # Set title
            title = slide.shapes.title
            title.text = slide_data.get('title', 'Untitled Slide')
            
            # Add content
            if len(slide.shapes) > 1 and hasattr(slide.shapes[1], 'text_frame'):
                body = slide.shapes[1]
                tf = body.text_frame
                tf.clear()  # Clear any existing content
                
                content_items = slide_data.get('content', [])
                for idx, item in enumerate(content_items):
                    if idx == 0:
                        tf.text = str(item)
                    else:
                        p = tf.add_paragraph()
                        p.text = str(item)
                        p.level = 0
            
            # Add speaker notes if present
            notes = slide_data.get('notes', '')
            if notes and hasattr(slide, 'notes_slide'):
                notes_slide = slide.notes_slide
                text_frame = notes_slide.notes_text_frame
                text_frame.text = notes
        
        # Save modified presentation
        output_path = f"{session_id}/presentation.pptx"
        with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as output_file:
            prs.save(output_file.name)
            
            # Upload to S3
            s3.upload_file(
                output_file.name,
                os.environ['ARTIFACTS_BUCKET_NAME'],
                output_path
            )
        
        # Generate presigned URL
        presigned_url = s3.generate_presigned_url(
            'get_object',
            Params={
                'Bucket': os.environ['ARTIFACTS_BUCKET_NAME'],
                'Key': output_path
            },
            ExpiresIn=3600
        )
        
        # Update session with document URL
        dynamodb.update_item(
            TableName=os.environ['SESSIONS_TABLE_NAME'],
            Key={'session_id': {'S': session_id}},
            UpdateExpression='SET document_urls = list_append(if_not_exists(document_urls, :empty), :url), #status = :status, updated_at = :ua',
            ExpressionAttributeNames={'#status': 'status'},
            ExpressionAttributeValues={
                ':url': {'L': [{'S': presigned_url}]},
                ':empty': {'L': []},
                ':status': {'S': 'POWERPOINT_GENERATED'},
                ':ua': {'S': datetime.utcnow().isoformat()}
            }
        )
        
        print(f"[POWERPOINT] PowerPoint generated successfully")
        
        return format_agent_response(200, {
            'session_id': session_id,
            'message': 'PowerPoint generated successfully',
            'document_url': presigned_url
        })
        
    except Exception as e:
        print(f"[POWERPOINT] ERROR: {str(e)}")
        import traceback
        print(f"[POWERPOINT] Traceback: {traceback.format_exc()}")
        return format_agent_response(500, f'PowerPoint generation failed: {str(e)}')

def format_agent_response(status_code, body):
    """Format response for Bedrock Agent action group"""
    if isinstance(body, dict):
        body_text = json.dumps(body)
    else:
        body_text = str(body)
    
    return {
        'messageVersion': '1.0',
        'response': {
            'actionGroup': 'PowerPointGenerator',
            'function': 'powerpointgenerator',  # Must match function name in action group
            'functionResponse': {
                'responseBody': {
                    'TEXT': {
                        'body': body_text
                    }
                }
            }
        }
    }
#!/usr/bin/env python3
import boto3
import os
import shutil
from botocore.exceptions import ClientError

SAMPLE_TEMPLATES = {
    'powerpoint': [
        {
            'name': 'basic_proposal.pptx',
            'slides': [
                {'title': 'Project Overview', 'layout': 0},
                {'title': 'Scope of Work', 'layout': 1},
                {'title': 'Timeline', 'layout': 1},
                {'title': 'Cost Breakdown', 'layout': 1},
                {'title': 'Team', 'layout': 1},
                {'title': 'Next Steps', 'layout': 1}
            ]
        }
    ],
    'sow': [
        {
            'name': 'standard_sow.docx',
            'sections': [
                'Executive Summary',
                'Project Overview',
                'Scope of Services',
                'Timeline and Milestones',
                'Deliverables',
                'Cost Structure',
                'Terms and Conditions'
            ]
        }
    ]
}

def get_bucket_name():
    cloudformation = boto3.client('cloudformation')
    try:
        response = cloudformation.describe_stacks(StackName='ScopeSmithInfrastructure')
        for output in response['Stacks'][0]['Outputs']:
            if output['OutputKey'] == 'TemplatesBucketName':
                return output['OutputValue']
    except ClientError as e:
        print(f"Error getting bucket name: {e}")
        return None

def create_sample_pptx(template):
    from pptx import Presentation
    
    prs = Presentation()
    for slide in template['slides']:
        layout = prs.slide_layouts[slide['layout']]
        slide_obj = prs.slides.add_slide(layout)
        
        if hasattr(slide_obj, 'shapes') and hasattr(slide_obj.shapes, 'title'):
            title = slide_obj.shapes.title
            if title:
                title.text = slide['title']
    
    temp_path = f"/tmp/{template['name']}"
    prs.save(temp_path)
    return temp_path

def create_sample_docx(template):
    from docx import Document
    
    doc = Document()
    
    # Add title
    doc.add_heading('Statement of Work', 0)
    
    # Add sections
    for section in template['sections']:
        doc.add_heading(section, level=1)
        doc.add_paragraph('Sample content for ' + section)
    
    temp_path = f"/tmp/{template['name']}"
    doc.save(temp_path)
    return temp_path

def upload_templates():
    bucket_name = get_bucket_name()
    if not bucket_name:
        print("Failed to get templates bucket name")
        return

    s3 = boto3.client('s3')
    
    # Create and upload PowerPoint templates
    for ppt_template in SAMPLE_TEMPLATES['powerpoint']:
        temp_path = create_sample_pptx(ppt_template)
        s3.upload_file(
            temp_path,
            bucket_name,
            f"powerpoint/{ppt_template['name']}"
        )
        os.remove(temp_path)
        print(f"Uploaded PowerPoint template: {ppt_template['name']}")
    
    # Create and upload SOW templates
    for sow_template in SAMPLE_TEMPLATES['sow']:
        temp_path = create_sample_docx(sow_template)
        s3.upload_file(
            temp_path,
            bucket_name,
            f"sow/{sow_template['name']}"
        )
        os.remove(temp_path)
        print(f"Uploaded SOW template: {sow_template['name']}")
    
    print("Successfully uploaded all sample templates")

if __name__ == "__main__":
    upload_templates()